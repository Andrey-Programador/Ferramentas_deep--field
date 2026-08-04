import io
import os
import re
import unicodedata
from datetime import datetime

import pandas as pd

AGENTES_POR_COMUNIDADE = {
    "SALGUEIRO": [
        "Anna Paula Santana", "Ana Paula Silva", "Ana Carolina",
        "Eliane de Paula", "Gislane da Cunha", "Josefina",
    ],
    "CORTE OITO": ["Nayara", "Andreza", "Nathalia", "Brenda", "Ícaro"],
    "CHACRINHA": ["Daiane", "Gabrielle", "Luciana", "Simone", "Tieres", "Anna Júlia"],
}

CORES_VERDES = ["#075E54", "#0B7368", "#0B8F7A", "#13A88E", "#35BDB1", "#68CEC5", "#9BE0D9"]


def normalizar(valor):
    texto = str(valor or "").strip().upper()
    texto = unicodedata.normalize("NFKD", texto)
    return " ".join("".join(c for c in texto if not unicodedata.combining(c)).split())


def comunidade_canonica(valor):
    """Converte variações encontradas no Excel para as três comunidades oficiais."""
    valor_n = normalizar(valor)
    compacta = re.sub(r"[^A-Z0-9]", "", valor_n)
    aliases = {
        "SALGUEIRO": {"SALGUEIRO", "COMUNIDADE SALGUEIRO", "ASRO SALGUEIRO"},
        "CORTE OITO": {"CORTE OITO", "CORTE 8", "CORTE8", "C8", "COMUNIDADE CORTE OITO", "ASRO CORTE OITO"},
        "CHACRINHA": {"CHACRINHA", "COMUNIDADE CHACRINHA", "ASRO CHACRINHA"},
    }
    for comunidade, nomes in aliases.items():
        normalizados = {normalizar(x) for x in nomes}
        compactos = {re.sub(r"[^A-Z0-9]", "", x) for x in normalizados}
        if valor_n in normalizados or compacta in compactos:
            return comunidade
        if any(nome in valor_n for nome in normalizados if len(nome) >= 6):
            return comunidade
    return "NAO RECONHECIDA"


def cadastro_agentes_df():
    linhas = []
    for comunidade, agentes in AGENTES_POR_COMUNIDADE.items():
        for agente in agentes:
            linhas.append({"ASRO": comunidade, "AGENTE": agente, "CHAVE_AGENTE": normalizar(agente), "STATUS_CADASTRO": "ATIVO"})
    return pd.DataFrame(linhas)


def _mapa_agentes():
    mapa = {}
    for comunidade, agentes in AGENTES_POR_COMUNIDADE.items():
        for agente in agentes:
            mapa[normalizar(agente)] = (agente, comunidade)
    return mapa


def agente_canonico(valor):
    """Reconhece diferenças de acento, caixa, espaços e nomes acompanhados da comunidade."""
    chave = normalizar(valor)
    mapa = _mapa_agentes()
    if chave in mapa:
        nome, comunidade = mapa[chave]
        return nome, normalizar(nome), comunidade, True

    # Aceita o nome cadastrado dentro de textos como 'Nayara - Corte Oito'.
    candidatos = []
    for chave_cadastrada, (nome, comunidade) in mapa.items():
        if chave_cadastrada in chave or chave in chave_cadastrada:
            # Evita associações inseguras de nomes muito curtos.
            if min(len(chave), len(chave_cadastrada)) >= 5:
                candidatos.append((nome, chave_cadastrada, comunidade))
    if len(candidatos) == 1:
        nome, chave_cadastrada, comunidade = candidatos[0]
        return nome, chave_cadastrada, comunidade, True
    return str(valor).strip() or "AGENTE NÃO INFORMADO", chave, None, False


def preparar_dados_relatorio(df_base):
    """Padroniza os dados filtrados e preserva registros não cadastrados."""
    base = df_base.copy()
    obrigatorias = ["ASRO", "AGENTE", "TIPO_CLASS"]
    faltantes = [c for c in obrigatorias if c not in base.columns]
    if faltantes:
        raise ValueError("O relatório não recebeu as colunas tratadas: " + ", ".join(faltantes))

    agente_info = base["AGENTE"].apply(agente_canonico)
    base["AGENTE_CANONICO"] = agente_info.map(lambda x: x[0])
    base["CHAVE_AGENTE"] = agente_info.map(lambda x: x[1])
    base["COMUNIDADE_PELO_AGENTE"] = agente_info.map(lambda x: x[2])
    base["AGENTE_CADASTRADO"] = agente_info.map(lambda x: x[3])
    base["COMUNIDADE_CANONICA"] = base["ASRO"].map(comunidade_canonica)

    # Se o ASRO vier como código ou texto desconhecido, usa a comunidade do agente cadastrado.
    usar_agente = (base["COMUNIDADE_CANONICA"] == "NAO RECONHECIDA") & base["COMUNIDADE_PELO_AGENTE"].notna()
    base.loc[usar_agente, "COMUNIDADE_CANONICA"] = base.loc[usar_agente, "COMUNIDADE_PELO_AGENTE"]

    base["TIPO_CLASS"] = base["TIPO_CLASS"].astype(str).map(normalizar)
    base["REGISTRO_RECONHECIDO"] = base["COMUNIDADE_CANONICA"].isin(AGENTES_POR_COMUNIDADE)
    return base


def diagnostico_relatorio(df_base):
    base = preparar_dados_relatorio(df_base)
    por_comunidade = {
        comunidade: int((base["COMUNIDADE_CANONICA"] == comunidade).sum())
        for comunidade in AGENTES_POR_COMUNIDADE
    }
    nao_reconhecidos = base[~base["REGISTRO_RECONHECIDO"]]
    agentes_nao_cadastrados = sorted(
        base.loc[base["REGISTRO_RECONHECIDO"] & ~base["AGENTE_CADASTRADO"], "AGENTE_CANONICO"]
        .dropna().astype(str).unique().tolist()
    )
    return {
        "total_filtrado": int(len(base)),
        "total_associado": int(base["REGISTRO_RECONHECIDO"].sum()),
        "total_nao_associado": int((~base["REGISTRO_RECONHECIDO"]).sum()),
        "por_comunidade": por_comunidade,
        "agentes_nao_cadastrados": agentes_nao_cadastrados,
        "asros_nao_reconhecidas": sorted(nao_reconhecidos["ASRO"].dropna().astype(str).unique().tolist()),
    }


def resumo_operacional_com_cadastro(df_base):
    """Inclui agentes cadastrados com zero e preserva agentes novos do Excel."""
    base = preparar_dados_relatorio(df_base)
    base = base[base["REGISTRO_RECONHECIDO"]].copy()

    producao = (
        base.groupby(["COMUNIDADE_CANONICA", "CHAVE_AGENTE", "AGENTE_CANONICO"], dropna=False)
        .agg(
            Visitas=("CHAVE_AGENTE", "size"),
            Adesoes=("TIPO_CLASS", lambda x: int((x == "ADESOES").sum())),
            Ausentes=("TIPO_CLASS", lambda x: int((x == "AUSENTES").sum())),
            Recusas=("TIPO_CLASS", lambda x: int((x == "RECUSAS").sum())),
            Agendamentos=("TIPO_CLASS", lambda x: int((x == "AGENDAMENTOS").sum())),
            Imoveis_vagos=("TIPO_CLASS", lambda x: int((x == "IMOVEIS VAGOS").sum())),
        ).reset_index()
        .rename(columns={"COMUNIDADE_CANONICA": "ASRO", "AGENTE_CANONICO": "AGENTE"})
    )

    cadastro = cadastro_agentes_df()[["ASRO", "AGENTE", "CHAVE_AGENTE"]].copy()
    resumo = cadastro.merge(
        producao.drop(columns=["AGENTE"]),
        on=["ASRO", "CHAVE_AGENTE"],
        how="left",
    )

    # Acrescenta agentes do Excel que ainda não constam no cadastro, sem perder visitas.
    chaves_cadastro = set(zip(cadastro["ASRO"], cadastro["CHAVE_AGENTE"]))
    extras = producao[
        ~producao.apply(lambda r: (r["ASRO"], r["CHAVE_AGENTE"]) in chaves_cadastro, axis=1)
    ].copy()
    if not extras.empty:
        resumo = pd.concat([resumo, extras], ignore_index=True, sort=False)

    metricas = ["Visitas", "Adesoes", "Ausentes", "Recusas", "Agendamentos", "Imoveis_vagos"]
    resumo[metricas] = resumo[metricas].fillna(0).astype(int)
    resumo["Taxa de adesão"] = resumo.apply(
        lambda r: r["Adesoes"] / r["Visitas"] if r["Visitas"] else 0, axis=1
    )
    resumo["Situação de sincronização"] = resumo["Visitas"].map(
        lambda v: "Com registros" if v > 0 else "Sem registros"
    )
    return resumo[["ASRO", "AGENTE", "Visitas", "Adesoes", "Ausentes", "Recusas", "Agendamentos", "Imoveis_vagos", "Taxa de adesão", "Situação de sincronização"]]


def _grafico_pizza_png(resumo_comunidade, titulo):
    import matplotlib.pyplot as plt
    dados = resumo_comunidade[resumo_comunidade["Visitas"] > 0].copy()
    fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=160)
    if dados.empty:
        ax.text(0.5, 0.5, "Nenhuma visita registrada", ha="center", va="center", fontsize=15, color="#075E54")
        ax.axis("off")
    else:
        labels = [f"{nome}\n{visitas} visitas" for nome, visitas in zip(dados["AGENTE"], dados["Visitas"])]
        cores = [CORES_VERDES[i % len(CORES_VERDES)] for i in range(len(dados))]
        ax.pie(
            dados["Visitas"], labels=labels, autopct=lambda p: f"{p:.1f}%".replace(".", ","),
            startangle=90, colors=cores, pctdistance=0.72,
            wedgeprops={"edgecolor": "white", "linewidth": 1.5},
            textprops={"fontsize": 8, "color": "#18312E"},
        )
        ax.axis("equal")
    ax.set_title(titulo, fontsize=14, fontweight="bold", color="#075E54", pad=12)
    fig.tight_layout()
    out = io.BytesIO()
    fig.savefig(out, format="png", bbox_inches="tight", transparent=True)
    plt.close(fig)
    out.seek(0)
    return out


def gerar_pdf_comunidade(df_base, comunidade, logo_path=None, periodo=""):
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, KeepTogether, PageBreak

    comunidade = comunidade_canonica(comunidade)
    resumo = resumo_operacional_com_cadastro(df_base)
    dados = resumo[resumo["ASRO"] == comunidade].copy().sort_values(["Visitas", "Adesoes", "AGENTE"], ascending=[False, False, True])
    if dados.empty:
        raise ValueError(f"Comunidade não cadastrada: {comunidade}")

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=14*mm, leftMargin=14*mm, topMargin=12*mm, bottomMargin=14*mm)
    styles = getSampleStyleSheet()
    titulo = ParagraphStyle("TituloSCA", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=20, leading=23, textColor=colors.white, alignment=TA_CENTER)
    subtitulo = ParagraphStyle("Subtitulo", parent=styles["Normal"], fontName="Helvetica", fontSize=10, textColor=colors.white, alignment=TA_CENTER)
    secao = ParagraphStyle("Secao", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=13, textColor=colors.HexColor("#075E54"), spaceBefore=8, spaceAfter=6)
    normal = ParagraphStyle("NormalBr", parent=styles["Normal"], fontName="Helvetica", fontSize=9, leading=12, textColor=colors.HexColor("#18312E"))

    total_visitas = int(dados["Visitas"].sum())
    total_adesoes = int(dados["Adesoes"].sum())
    com_registros = int((dados["Visitas"] > 0).sum())
    sem_registros = int((dados["Visitas"] == 0).sum())
    taxa = total_adesoes / total_visitas if total_visitas else 0
    periodo = periodo or "Período selecionado no sistema"

    logo = None
    if logo_path and os.path.exists(logo_path):
        logo = Image(logo_path, width=42*mm, height=8.8*mm)
    cab_esq = logo if logo else Paragraph("<b>Light</b>", titulo)
    cab = Table([[cab_esq, Paragraph("RELATÓRIO DE PRODUTIVIDADE", titulo)], ["", Paragraph(f"{comunidade.title()} | {periodo}", subtitulo)]], colWidths=[48*mm, 125*mm])
    cab.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#0B8F7A")), ("SPAN", (0,0), (0,1)),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"), ("ALIGN", (1,0), (1,-1), "CENTER"),
        ("LEFTPADDING", (0,0), (-1,-1), 10), ("RIGHTPADDING", (0,0), (-1,-1), 10),
        ("TOPPADDING", (0,0), (-1,-1), 10), ("BOTTOMPADDING", (0,0), (-1,-1), 10),
    ]))

    cards = [
        ("Total de visitas", total_visitas), ("Total de adesões", total_adesoes),
        ("Agentes cadastrados", len(dados)), ("Com registros", com_registros),
        ("Sem registros", sem_registros), ("Taxa de adesão", f"{taxa:.1%}".replace(".", ",")),
    ]
    rows = []
    for i in range(0, 6, 3):
        rows.append([Paragraph(f"<b>{cards[j][1]}</b><br/><font size='8'>{cards[j][0]}</font>", normal) for j in range(i, i+3)])
    cards_tbl = Table(rows, colWidths=[57*mm]*3, rowHeights=[20*mm, 20*mm], hAlign="CENTER")
    cards_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#E5F6F2")),
        ("BOX", (0,0), (-1,-1), 0.7, colors.HexColor("#13A88E")),
        ("INNERGRID", (0,0), (-1,-1), 4, colors.white),
        ("ALIGN", (0,0), (-1,-1), "CENTER"), ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
    ]))

    pizza = _grafico_pizza_png(dados, f"Distribuição das visitas por agente - {comunidade.title()}")
    pizza_img = Image(pizza, width=165*mm, height=96*mm)
    sem_lista = ", ".join(dados.loc[dados["Visitas"] == 0, "AGENTE"].tolist()) or "Nenhum"

    tabela_dados = [["Agente", "Visitas", "Adesões", "Ausentes", "Recusas", "Agend.", "Taxa"]]
    for _, r in dados.iterrows():
        tabela_dados.append([
            r["AGENTE"], int(r["Visitas"]), int(r["Adesoes"]), int(r["Ausentes"]), int(r["Recusas"]),
            int(r["Agendamentos"]), f"{r['Taxa de adesão']:.1%}".replace(".", ","),
        ])
    tabela = Table(tabela_dados, colWidths=[51*mm, 18*mm, 18*mm, 19*mm, 18*mm, 18*mm, 20*mm], repeatRows=1)
    tabela.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#075E54")), ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"), ("ALIGN", (1,1), (-1,-1), "CENTER"),
        ("FONTNAME", (0,1), (-1,-1), "Helvetica"), ("FONTSIZE", (0,0), (-1,-1), 8),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#E5F6F2")]),
        ("GRID", (0,0), (-1,-1), 0.4, colors.HexColor("#68CEC5")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"), ("TOPPADDING", (0,0), (-1,-1), 5), ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))

    story = [cab, Spacer(1, 8*mm), cards_tbl, Spacer(1, 6*mm), Paragraph("Produtividade por agente", secao), pizza_img,
             Paragraph(f"<b>Agentes sem registros no período:</b> {sem_lista}", normal), Spacer(1, 5*mm),
             Paragraph("Detalhamento individual", secao), tabela]

    def rodape(canvas, _doc):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#13A88E")); canvas.line(14*mm, 11*mm, 196*mm, 11*mm)
        canvas.setFillColor(colors.HexColor("#667085")); canvas.setFont("Helvetica", 7.5)
        canvas.drawString(14*mm, 7*mm, "Deep Field | Automação e Inteligência Operacional")
        canvas.drawRightString(196*mm, 7*mm, f"Página {_doc.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=rodape, onLaterPages=rodape)
    buffer.seek(0)
    return buffer.getvalue()


def gerar_pdf_mensal(df_base, logo_path=None, periodo="Relatório mensal"):
    from pypdf import PdfReader, PdfWriter
    writer = PdfWriter()
    for comunidade in ["SALGUEIRO", "CORTE OITO", "CHACRINHA"]:
        pdf = gerar_pdf_comunidade(df_base, comunidade, logo_path, periodo)
        reader = PdfReader(io.BytesIO(pdf))
        for page in reader.pages:
            writer.add_page(page)
    out = io.BytesIO(); writer.write(out); out.seek(0)
    return out.getvalue()


def adicionar_controle_agentes_excel(excel_bytes, df_base):
    from openpyxl import load_workbook
    from openpyxl.styles import PatternFill, Font, Alignment
    wb = load_workbook(io.BytesIO(excel_bytes))
    if "CONTROLE DE AGENTES" in wb.sheetnames:
        del wb["CONTROLE DE AGENTES"]
    ws = wb.create_sheet("CONTROLE DE AGENTES", 1)
    resumo = resumo_operacional_com_cadastro(df_base)
    headers = ["COMUNIDADE", "AGENTE", "VISITAS", "ADESÕES", "AUSENTES", "RECUSAS", "AGENDAMENTOS", "IMÓVEIS VAGOS", "TAXA DE ADESÃO", "SINCRONIZAÇÃO"]
    ws.append(headers)
    for _, r in resumo.iterrows():
        ws.append([r["ASRO"], r["AGENTE"], r["Visitas"], r["Adesoes"], r["Ausentes"], r["Recusas"], r["Agendamentos"], r["Imoveis_vagos"], r["Taxa de adesão"], r["Situação de sincronização"]])
    for c in ws[1]:
        c.fill = PatternFill("solid", fgColor="075E54"); c.font = Font(color="FFFFFF", bold=True); c.alignment = Alignment(horizontal="center")
    for row in range(2, ws.max_row + 1):
        ws.cell(row, 9).number_format = "0.0%"
        if ws.cell(row, 10).value == "Sem registros":
            for c in ws[row]: c.fill = PatternFill("solid", fgColor="FFF3CD")
    widths = [18, 28, 12, 12, 12, 12, 16, 18, 18, 22]
    for idx, width in enumerate(widths, 1): ws.column_dimensions[chr(64+idx)].width = width
    ws.freeze_panes = "A2"; ws.auto_filter.ref = ws.dimensions
    out = io.BytesIO(); wb.save(out); out.seek(0)
    return out.getvalue()


def gerar_powerpoint_resumo(df_base, logo_path=None, periodo="Período selecionado"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    resumo = resumo_operacional_com_cadastro(df_base)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    verde = RGBColor(19,168,142); escuro = RGBColor(7,94,84); branco = RGBColor(255,255,255); claro = RGBColor(229,246,242)

    def add_header(slide, titulo, subtitulo=""):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.75)); shape.fill.solid(); shape.fill.fore_color.rgb = verde; shape.line.fill.background()
        if logo_path and os.path.exists(logo_path): slide.shapes.add_picture(logo_path, Inches(.35), Inches(.18), width=Inches(1.8))
        tb = slide.shapes.add_textbox(Inches(2.4), Inches(.12), Inches(8.4), Inches(.45)); p=tb.text_frame.paragraphs[0]; p.text=titulo; p.font.size=Pt(22); p.font.bold=True; p.font.color.rgb=branco; p.alignment=PP_ALIGN.CENTER
        if subtitulo:
            st=slide.shapes.add_textbox(Inches(10.7), Inches(.18), Inches(2.2), Inches(.3)); p=st.text_frame.paragraphs[0]; p.text=subtitulo; p.font.size=Pt(9); p.font.color.rgb=branco; p.alignment=PP_ALIGN.RIGHT

    slide = prs.slides.add_slide(blank); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=verde
    if logo_path and os.path.exists(logo_path): slide.shapes.add_picture(logo_path, Inches(4.25), Inches(1.25), width=Inches(4.8))
    tb=slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(1.4)); tf=tb.text_frame; p=tf.paragraphs[0]; p.text="Deep Field"; p.font.size=Pt(32); p.font.bold=True; p.font.color.rgb=branco; p.alignment=PP_ALIGN.CENTER
    p=tf.add_paragraph(); p.text="Automação e Inteligência Operacional"; p.font.size=Pt(20); p.font.color.rgb=branco; p.alignment=PP_ALIGN.CENTER
    p=tf.add_paragraph(); p.text=periodo; p.font.size=Pt(13); p.font.color.rgb=branco; p.alignment=PP_ALIGN.CENTER

    for comunidade in ["SALGUEIRO", "CORTE OITO", "CHACRINHA"]:
        dados=resumo[resumo["ASRO"]==comunidade].copy(); slide=prs.slides.add_slide(blank); add_header(slide, f"RESULTADOS | {comunidade}", periodo)
        total_v=int(dados["Visitas"].sum()); total_a=int(dados["Adesoes"].sum()); ativos=int((dados["Visitas"]>0).sum()); sem=int((dados["Visitas"]==0).sum())
        cards=[("Visitas", total_v), ("Adesões", total_a), ("Com registros", ativos), ("Sem registros", sem)]
        for i,(rot,val) in enumerate(cards):
            x=.55+i*3.15; sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.05), Inches(2.75), Inches(1.05)); sh.fill.solid(); sh.fill.fore_color.rgb=claro; sh.line.color.rgb=verde
            tf=sh.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=str(val); p.font.size=Pt(24); p.font.bold=True; p.font.color.rgb=escuro; p.alignment=PP_ALIGN.CENTER; p=tf.add_paragraph(); p.text=rot; p.font.size=Pt(11); p.font.color.rgb=escuro; p.alignment=PP_ALIGN.CENTER
        prod=dados[dados["Visitas"]>0]
        chart_data=ChartData(); chart_data.categories=prod["AGENTE"].tolist() or ["Sem registros"]; chart_data.add_series("Visitas", prod["Visitas"].tolist() or [1])
        chart=slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(.55), Inches(2.35), Inches(6.15), Inches(4.55), chart_data).chart; chart.has_legend=True; chart.legend.position=XL_LEGEND_POSITION.RIGHT; chart.legend.font.size=Pt(10); chart.has_title=True; chart.chart_title.text_frame.text="Visitas por agente"
        rows=len(dados)+1; cols=4; table=slide.shapes.add_table(rows, cols, Inches(7.0), Inches(2.35), Inches(5.75), Inches(4.55)).table
        for j,h in enumerate(["Agente","Visitas","Adesões","Status"]): table.cell(0,j).text=h
        for i,(_,r) in enumerate(dados.sort_values("Visitas",ascending=False).iterrows(),1):
            for j,v in enumerate([r["AGENTE"],str(r["Visitas"]),str(r["Adesoes"]),r["Situação de sincronização"]]): table.cell(i,j).text=str(v)
        for r in range(rows):
            for c in range(cols):
                cell=table.cell(r,c); cell.text_frame.paragraphs[0].font.size=Pt(9); cell.text_frame.paragraphs[0].font.color.rgb = branco if r==0 else escuro; cell.fill.solid(); cell.fill.fore_color.rgb = escuro if r==0 else (claro if r%2==0 else branco)
    out=io.BytesIO(); prs.save(out); out.seek(0); return out.getvalue()
